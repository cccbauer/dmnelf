#!/usr/bin/env python3
"""
Build a summary PPTX for the mindwear EEG->BOLD decoder investigation: why the ball demo doesn't
visually "line up," what that traced back to (an in-sample-only accuracy claim), and everything
tried to improve it (all negative). Companion to build_ppt.py's EFP replication deck -- same style.

Source numbers not available in this repo (Phase 1/2/2.5 tables) come from the `efp-pooled` branch
(analysis/fingerprint/efp_pooled/HANDOFF.md and results/*_summary.json there) -- hardcoded here
with inline citations since that branch isn't merged to main. The Phase 3 calibration numbers ARE
local (results/efp_calibrated_loso_epoc12_multi_summary.csv, committed to main).
"""
from pathlib import Path
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PROJ = Path(__file__).resolve().parent.parent
RES = PROJ / "results"
OUT_FIGS = RES / "decoder_investigation_figs"
OUT_FIGS.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x1F, 0x3A, 0x5F); GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32); RED = RGBColor(0xB0, 0x3A, 0x2E)
C_NAVY, C_GREEN, C_RED, C_GREY = "#1F3A5F", "#2E7D32", "#B03A2E", "#555555"

# ── Fig 1: Phase 1 -- in-sample vs honest held-out (efp-pooled branch, shipped_on_*_summary.json,
#    commit e308996/57fc330; see HANDOFF.md "PHASE 1 RESULT") ──
targets = ["CEN", "DMN", "PDA"]
dmnelf_insample = {"CEN": 0.171, "DMN": 0.249, "PDA": 0.194}
rtbpd_nf1_heldout = {"CEN": 0.069, "DMN": 0.052, "PDA": 0.010}
rtbpd_nf1_p = {"CEN": 0.0004, "DMN": 0.003, "PDA": 0.43}
rtbpd_nf2_heldout = {"CEN": 0.093, "DMN": 0.036, "PDA": 0.008}
rtbpd_nf2_p = {"CEN": 0.006, "DMN": 0.11, "PDA": 0.71}

fig, ax = plt.subplots(figsize=(9, 4.3))
x = np.arange(len(targets)); w = 0.26
bars1 = ax.bar(x - w, [dmnelf_insample[t] for t in targets], w, label="DMNELF (in-sample)", color=C_GREY, edgecolor="k", linewidth=0.5)
bars2 = ax.bar(x, [rtbpd_nf1_heldout[t] for t in targets], w, label="rtBPD nf1 (held out)", color=C_NAVY, edgecolor="k", linewidth=0.5)
bars3 = ax.bar(x + w, [rtbpd_nf2_heldout[t] for t in targets], w, label="rtBPD nf2 (held out)", color="#4C8FC9", edgecolor="k", linewidth=0.5)
for bars, ps in ((bars2, rtbpd_nf1_p), (bars3, rtbpd_nf2_p)):
    for b, t in zip(bars, targets):
        sig = "*" if ps[t] < 0.05 else "n.s."
        ax.annotate(sig, (b.get_x() + b.get_width() / 2, b.get_height()), textcoords="offset points",
                    xytext=(0, 3), ha="center", fontsize=10,
                    color=(C_GREEN if sig == "*" else C_RED), fontweight="bold")
ax.axhline(0, color="k", lw=0.8); ax.set_xticks(x); ax.set_xticklabels(targets)
ax.set_ylabel("mean r"); ax.set_title("The shipped decoder's PDA has no out-of-sample validity")
ax.legend(fontsize=10); fig.tight_layout()
fig1 = OUT_FIGS / "fig1_insample_vs_heldout.png"; fig.savefig(fig1, dpi=140); plt.close(fig)

# ── Fig 2: Phase 2 + 2.5 -- grouped-CV comparison across all arms tried (efp-pooled branch,
#    HANDOFF.md "PHASE 2 FIRST PASS" + "PHASE 2.5" sections, commits cdf1f92/340eda5) ──
arms = ["baseline\n(shipped arch.)", "pooled28\n(+9 rtBPD)", "epoc_afproxy\n(+Fp1/Fp2)",
        "cap31\n(full montage)", "ElasticNet", "joint PLS\n(CEN+DMN)"]
cen_vals = [0.042, 0.075, 0.041, 0.049, 0.044, 0.012]
dmn_vals = [0.032, 0.037, 0.030, 0.019, 0.019, 0.012]

fig, ax = plt.subplots(figsize=(10.5, 4.3))
x = np.arange(len(arms)); w = 0.35
ax.bar(x - w / 2, cen_vals, w, label="CEN", color=C_NAVY, edgecolor="k", linewidth=0.5)
ax.bar(x + w / 2, dmn_vals, w, label="DMN", color="#4C8FC9", edgecolor="k", linewidth=0.5)
ax.axhline(cen_vals[0], color=C_NAVY, lw=1, ls="--", alpha=0.5)
ax.axhline(dmn_vals[0], color="#4C8FC9", lw=1, ls="--", alpha=0.5)
ax.axhline(0, color="k", lw=0.8)
ax.set_xticks(x); ax.set_xticklabels(arms, fontsize=9)
ax.set_ylabel("grouped-CV r (training-side, not held-out)")
ax.set_title("Every structural change tried: none clearly beats the shipped architecture")
ax.legend(fontsize=10); fig.tight_layout()
fig2 = OUT_FIGS / "fig2_arms_comparison.png"; fig.savefig(fig2, dpi=140); plt.close(fig)

# ── Fig 3: Phase 3 -- per-subject calibration, local CSV (results/efp_calibrated_loso_epoc12_multi_summary.csv) ──
cal = pd.read_csv(RES / "efp_calibrated_loso_epoc12_multi_summary.csv")

fig, axes = plt.subplots(1, 2, figsize=(10.5, 4.0), sharey=True)
for ax, target, color in zip(axes, ["CEN", "DMN"], [C_NAVY, "#4C8FC9"]):
    d = cal[cal.target == target]
    n_cals = [1, 2, 3]   # n_cal=0 has no pseudo_cal (no calibration data yet) -- skip it here
    g = [d[(d.n_cal_runs == n) & (d.method == "group_only")]["mean_r"].values[0] for n in n_cals]
    p = [d[(d.n_cal_runs == n) & (d.method == "pseudo_cal")]["mean_r"].values[0] for n in n_cals]
    ax.plot(n_cals, g, "o-", color=C_GREY, lw=2, label="group_only (baseline)")
    ax.plot(n_cals, p, "o-", color=(C_RED if target == "DMN" else color), lw=2, label="pseudo_cal (calibrated)")
    ax.axhline(0, color="k", lw=0.6)
    ax.set_title(target, fontweight="bold"); ax.set_xlabel("# calibration runs used")
    ax.set_xticks(n_cals)
    ax.legend(fontsize=9)
axes[0].set_ylabel("mean r (honest cross-subject LOSO, n=19)")
fig.suptitle("Per-subject calibration: no help for CEN, systematically worse for DMN", fontweight="bold")
fig.tight_layout()
fig3 = OUT_FIGS / "fig3_calibration.png"; fig.savefig(fig3, dpi=140); plt.close(fig)

# ── Presentation ──
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def title(s, text, sub=None):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)); tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p = tf.add_paragraph(); rr = p.add_run(); rr.text = sub; rr.font.size = Pt(15); rr.font.color.rgb = GREY


def bullets(s, items, top=1.5, left=0.6, width=12.1, height=5.5, size=17):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)); tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(items):
        lvl = 0; txt = b; col = None; bold = False
        if isinstance(b, tuple):
            txt, lvl = b[0], b[1]
            col = b[2] if len(b) > 2 else None; bold = b[3] if len(b) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl)
        if col: r.font.color.rgb = col
        r.font.bold = bold; p.space_after = Pt(6)


def img(s, path, left, top, width=None, height=None):
    if Path(path).exists():
        kw = {}
        if width: kw["width"] = Inches(width)
        if height: kw["height"] = Inches(height)
        s.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


# 1 title
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.5)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "The mindwear EEG→BOLD decoder: what actually transfers"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
for line, sz in [("Auditing the deployed decoder's real generalization, and everything tried to improve it", 18),
                 ("efp_pooled investigation, Aug 31 – Sep 1 2026", 14)]:
    p = tf.add_paragraph(); rr = p.add_run(); rr.text = line; rr.font.size = Pt(sz); rr.font.color.rgb = GREY

# 2 motivation
s = prs.slides.add_slide(BLANK); title(s, "Why we looked at this")
bullets(s, [
    "The Compare-mode demo shows two balls side by side — one driven by observed fMRI BOLD, one "
    "by the portable EEG decoder — meant to visually track each other.",
    "They don't line up well in practice. Tried display-only fixes first (EMA smoothing, "
    "outlier winsorizing/median filtering) — reverted or inconclusive.",
    ("Root cause was NOT display noise: simulating the real ball physics against real BOLD data "
     "showed the visible \"whips\" line up with genuine, large swings in the underlying signal.", 0, NAVY, True),
    "That redirected the question from \"how do we smooth the display\" to \"how accurate is the "
    "decoder, really\" — which had never been honestly, out-of-sample tested before.",
], size=17)

# 3 Phase 1
s = prs.slides.add_slide(BLANK); title(s, "The quoted accuracy number was in-sample")
img(s, fig1, left=1.5, top=1.5, width=10.3)
bullets(s, [
    ("The r≈+0.22 quoted for the decoder was scored on dmnelf005 — a subject IN the training "
     "set, and the cohort's single best performer. Never scored on a truly new subject before.", 0),
    ("Scored honestly on rtBPD subjects (never in training): PDA r=+0.010 (p=0.43) — a coin flip. "
     "CEN and DMN individually transfer significantly; their difference (PDA) does not.", 0, RED, True),
], top=6.2, size=14)

# 4 why PDA fails / decision
s = prs.slides.add_slide(BLANK); title(s, "Why PDA specifically fails, and what we kept")
bullets(s, [
    "The EEG-decodable component is largely shared/global (consistent with an arousal-type "
    "signal, see eeg_bold_coupling analysis) — real, but not network-specific.",
    ("PDA = CEN − DMN differences away exactly that shared component, leaving nothing "
     "transferable behind. Fitting PDA directly instead of deriving it (obvious fix) does not "
     "recover it either — tested and negative.", 0, RED, True),
    ("Decision: keep the shipped two-model architecture (CEN and DMN each fit independently, PDA "
     "= their difference downstream), and focus improvement efforts on CEN/DMN — the targets "
     "that actually transfer.", 0, NAVY, True),
], size=17)

# 5 arms tried
s = prs.slides.add_slide(BLANK); title(s, "What we tried to improve CEN/DMN transfer")
img(s, fig2, left=1.2, top=1.5, width=11.0)
bullets(s, [
    ("pooled28's CEN bar looks like a clear win here (this is training-side grouped-CV) — it "
     "did NOT hold up on the locked external test (CEN p worsened 0.038→0.117). Grouped-CV "
     "overstated transfer; the chart above is a training diagnostic, not a performance claim.", 0),
    ("cap31 (full research montage) is the only other real effect — but it trades a small CEN "
     "gain for a real DMN loss, not a net win. epoc_afproxy and ElasticNet are flat-to-worse.", 0),
    ("The joint CEN+DMN fit (motivated by the shared-component finding) is the WORST arm despite "
     "having the HIGHEST in-sample r — a clean overfitting signature.", 0, RED, True),
], top=6.3, size=13)

# 6 calibration
s = prs.slides.add_slide(BLANK); title(s, "Per-subject calibration — the most promising lever — also negative")
img(s, fig3, left=1.2, top=1.45, width=11.0)
bullets(s, [
    ("Calibrate on the subject's own run-1 task design (no fMRI needed — the deployable "
     "scenario), test on their remaining real fMRI, honest cross-subject LOSO (n=19).", 0),
    ("CEN never significantly improves; DMN gets systematically WORSE, and more so with more "
     "calibration data. (PDA/GSR_PDA do improve — but that's not the kept target.)", 0, RED, True),
], top=6.35, size=14)

# 7 conclusion
s = prs.slides.add_slide(BLANK); title(s, "Conclusion")
bullets(s, [
    ("Six independent approaches tried — pooling, direct-PDA-fit, wider montage, ElasticNet, "
     "joint multi-task fitting, per-subject calibration — none improves CEN/DMN transfer beyond "
     "the currently-shipped model.", 0, NAVY, True),
    "The ceiling looks set by how much EEG-decodable signal exists at all (largely shared/"
    "global), not by any modeling choice tried so far.",
    "Recommendation: keep the shipped efp_epoc_model.npz as-is. If a single, most-defensible "
    "target is needed for messaging, DMN transfers most consistently across every arm tested.",
    ("Expectation-setting for demos: dual-ball direction agreement needs r≈0.59 for 70%; the "
     "validated ceiling here is r≈0.15–0.25. A visually \"convincing\" side-by-side demo isn't "
     "reachable by any decoding fix on the table — this was never really a display problem.", 0, GREY),
], size=17)

out = PROJ / "efp_decoder_investigation_results.pptx"; prs.save(str(out))
print(f"Saved {out} ({len(prs.slides._sldIdLst)} slides)")
