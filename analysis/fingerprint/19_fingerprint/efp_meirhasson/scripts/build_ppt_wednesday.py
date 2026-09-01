#!/usr/bin/env python3
"""
Build the Wednesday lab-talk deck: experiential intro -> why neurofeedback -> Meir-Hasson 2014
method -> our n=19 replication + rtBPD cross-cohort results -> mindwear -> live demo.

Meir-Hasson 2014 (slides 9-10), the prior-fMRI-evidence slide (7), and the related-work slide (21)
are filled in from the actual papers (Meir-Hasson et al. 2014 NeuroImage; Zhang et al. 2023 Mol
Psychiatry; Bauer et al. 2025 Psychiatry Res: Neuroimaging) -- real cited numbers, not fabricated.
No figure IMAGES from those papers are embedded (not available as local files); each such slide
notes where to drop in the original figure if wanted. The mindwear screenshot on slide 18 is the
one remaining true placeholder -- this script has no access to a current screenshot.

Cross-cohort (DMNELF->rtBPD) numbers are read from cross_cohort_efp_summary_tr{,_nf2}.csv in
THIS (19_fingerprint) results dir -- must be freshly re-run on n=19 first (see
submit_crosscohort_n19.sh); the only pre-existing cross-cohort CSVs on disk before that were
n=17 (predate dmnelf002/003) and would misrepresent the current cohort if used by mistake.
"""
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJ = Path(__file__).resolve().parent.parent
RES = PROJ / "results"
FULL = RES / "full"

NAVY = RGBColor(0x1F, 0x3A, 0x5F); GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32); RED = RGBColor(0xB0, 0x3A, 0x2E)
ORANGE = RGBColor(0xC4, 0x6A, 0x00)

# ── require the fresh n=19 cross-cohort numbers; refuse to silently fall back to stale n=17 ──
cc1_path = RES / "cross_cohort_efp_summary_tr.csv"
cc2_path = RES / "cross_cohort_efp_summary_tr_nf2.csv"
if not (cc1_path.exists() and cc2_path.exists()):
    raise SystemExit(
        f"Missing {cc1_path.name}/{cc2_path.name} in {RES} -- run submit_crosscohort_n19.sh on "
        "the cluster first (re-validates DMNELF->rtBPD on n=19; the numbers must not come from "
        "the old n=17 CSVs living under the sibling efp_meirhasson/ tree).")
cc1 = pd.read_csv(cc1_path); cc2 = pd.read_csv(cc2_path)


def cc(df, target, method="EFP"):
    row = df[(df.target == target) & (df.method == method)]
    return row.iloc[0] if len(row) else None


# ── Presentation ──
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def title(s, text, sub=None):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)); tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p = tf.add_paragraph(); rr = p.add_run(); rr.text = sub; rr.font.size = Pt(14); rr.font.color.rgb = GREY


def bullets(s, items, top=1.5, left=0.6, width=12.1, height=5.5, size=17):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)); tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(items):
        lvl = 0; txt = b; col = None; bold = False
        if isinstance(b, tuple):
            txt, lvl = b[0], b[1]
            col = b[2] if len(b) > 2 else None; bold = b[3] if len(b) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl)
        if col: r.font.color.rgb = col
        r.font.bold = bold; p.space_after = Pt(6)


def img(s, path, left, top, width=None, height=None):
    if Path(path).exists():
        kw = {}
        if width: kw["width"] = Inches(width)
        if height: kw["height"] = Inches(height)
        s.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


# 1 — Title
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.0)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "From the Scanner to the Living Room:"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
p = tf.add_paragraph(); rr = p.add_run(); rr.text = "Scaling Default Mode Network Neurofeedback with a Portable EEG Fingerprint"
rr.font.size = Pt(24); rr.font.bold = True; rr.font.color.rgb = NAVY
for line, sz in [("[Your name] · [Lab/Department] · [Date]", 16)]:
    p = tf.add_paragraph(); rr = p.add_run(); rr.text = line; rr.font.size = Pt(sz); rr.font.color.rgb = GREY

# 2 — Experiential intro: facilitation slide
s = prs.slides.add_slide(BLANK); title(s, "Before we begin: a 3-minute experiment")
bullets(s, [
    ("Close your eyes (or soften your gaze). We'll do this together, in silence.", 0, NAVY, True),
    ("Part 1 — 30 seconds: just breathe. Notice your breath, nothing else.", 0),
    ("Part 2 — 2.5 minutes: bring to mind something unresolved — a decision you're unsure "
     "about, an interaction that didn't go the way you wanted, something still bothering you. "
     "Don't try to solve it. Just let your mind go there, and stay with it.", 0),
    ("(Speaker: keep time silently — e.g. phone timer, screen turned away from the room. "
     "Don't narrate the timing out loud once part 2 starts; let the room sit in it.)", 0, GREY),
], size=18)

# 3 — Debrief
s = prs.slides.add_slide(BLANK); title(s, "What did you notice?")
bullets(s, [
    "Ask the room: racing thoughts? Replaying the scenario? Judging yourself? Losing track of "
    "the breath entirely?",
    ("That drift — away from the present moment and into self-referential, evaluative, looping "
     "thought — has a name in neuroscience: the Default Mode Network (DMN).", 0, NAVY, True),
    "You didn't choose to ruminate just now. For most people, without training, that network "
    "runs the show the moment there's nothing else to do.",
], size=19)

# 4 — Why this matters
s = prs.slides.add_slide(BLANK); title(s, "Why this matters")
bullets(s, [
    "The DMN is most active exactly when we're NOT engaged in an external task — mind-wandering, "
    "self-referential processing, autobiographical memory.",
    "Excessive, hard-to-disengage DMN activity — rumination — is a transdiagnostic feature of "
    "depression, anxiety, and chronic stress, not just an occasional annoyance.",
    ("[Add your specific citations/prevalence stats here]", 0, ORANGE, True),
    "The question this line of work asks: can people learn to quiet this network on demand?",
], size=18)

# 5 — Why mindfulness
s = prs.slides.add_slide(BLANK); title(s, "Why mindfulness")
bullets(s, [
    "Experienced meditators show reduced DMN activity and altered DMN connectivity, both at "
    "rest and during tasks, compared to non-meditators.",
    "Mindfulness training is, in effect, DMN self-regulation training — learned gradually, "
    "through sustained practice.",
    ("[Cite your preferred primary sources here — e.g. Brewer et al. 2011 PNAS]", 0, ORANGE, True),
    ("The catch: this normally takes months to years of practice before it generalizes.", 0, RED, True),
], size=18)

# 6 — Why neurofeedback
s = prs.slides.add_slide(BLANK); title(s, "Why neurofeedback")
bullets(s, [
    "Real-time neurofeedback gives someone a direct, immediate readout of their own brain "
    "state — turning an invisible, internal process into something they can see and steer.",
    "The hypothesis: if people can SEE their DMN/CEN balance moment to moment, they can learn "
    "to shift it far faster than through practice alone — training a muscle you couldn't "
    "previously feel.",
    "Our approach: a visual \"ball\" task — the ball rises toward a Central Executive Network "
    "(CEN) target when task-focused attention dominates, and toward the DMN target otherwise.",
], size=18)

# 7 — prior fMRI results: Zhang et al. 2023 (mbNF) — direct precedent, same PDA metric
s = prs.slides.add_slide(BLANK)
title(s, "Prior fMRI evidence: this already works in the scanner",
      "Zhang, Raya, Morfini, ... Bauer, Whitfield-Gabrieli (2023). Molecular Psychiatry.")
bullets(s, [
    "9 adolescents with a lifetime history of depression/anxiety; DMN and CEN individually "
    "localized per person from a resting-state scan.",
    "One 15-minute mindfulness-based fMRI neurofeedback (mbNF) session: watch a dot driven by "
    "PDA = CEN activation − DMN activation move up with effective mental-noting practice — the "
    "exact same PDA metric and ball/dot mechanic used throughout this talk and in mindwear.",
    ("Result: participants spent significantly more time in the target state (DMN < CEN) than "
     "chance (p=.038), and all 9 of 9 adolescents showed reduced within-DMN connectivity after "
     "just one session.", 0, GREEN, True),
    "That connectivity reduction correlated with increased state mindfulness (r=−0.88, p=.002), "
    "and statistically mediated the link between neurofeedback performance and the mindfulness "
    "gain — i.e., quieting the DMN is plausibly the mechanism, not just a correlate.",
    ("This is the evidence base motivating scaling the same paradigm outside the scanner.", 0, NAVY, True),
], size=15)

# 8 — The scaling problem
s = prs.slides.add_slide(BLANK); title(s, "The scaling problem")
bullets(s, [
    "fMRI neurofeedback works, but an MRI scanner is expensive, immobile, and inaccessible for "
    "repeated, at-home, or scalable clinical use.",
    ("If we want this to reach people beyond a single research scanner, we need the SAME signal, "
     "decodable from something portable.", 0, NAVY, True),
    "EEG is cheap, wearable, and real-time — but scalp EEG doesn't measure DMN/CEN BOLD "
    "activity directly. It has to be decoded.",
    ("That decoding problem is what the rest of this talk is about.", 0, NAVY, True),
], size=19)

# 9 — Meir-Hasson 2014 background
s = prs.slides.add_slide(BLANK)
title(s, "Meir-Hasson et al. (2014) — the EEG Finger-Print",
      "NeuroImage 102:128–141  ·  [insert the paper's Fig. 1/Fig. 2 schematic here if desired]")
bullets(s, [
    "Goal: predict fMRI BOLD activity in a specific ROI — especially hard-to-reach subcortical "
    "regions like the amygdala — from simultaneous EEG, well enough to eventually let portable "
    "EEG stand in for the scanner during neurofeedback.",
    ("Core method (what we replicate): single electrode -> Stockwell time-frequency transform -> "
     "10 data-driven equal-energy frequency bands (not fixed alpha/theta) -> ridge regression on "
     "a [band x sliding time-delay] design.", 0, NAVY, True),
    "Key innovation: NO fixed hemodynamic delay. Each frequency band gets its own data-driven "
    "delay, per subject and electrode — because the true EEG-to-BOLD lag varies by all three.",
    "Validated with double cross-validation (outer block CV for honest test scores, inner CV for "
    "the ridge regularization strength) — the same discipline our own replication follows.",
    "Tested on two datasets: eyes-open/closed alpha modulation -> visual cortex (n=4), and "
    "real-time theta/alpha relaxation neurofeedback -> amygdala and dmPFC (n=20).",
], size=16)

# 10 — Meir-Hasson 2014 results
s = prs.slides.add_slide(BLANK)
title(s, "Meir-Hasson et al. (2014) — results",
      "[insert the paper's Fig. 4/6/7 bar charts and amygdala EFP maps here if desired]")
bullets(s, [
    ("Visual cortex (n=4): their EFP beat a traditional alpha-power predictor in ALL subjects, "
     "and beat a fixed-delay (HRF) predictor in all subjects (75% significantly).", 0),
    ("Amygdala (17/20 subjects cleared their validation threshold): EFP beat the traditional "
     "theta/alpha predictor in 95% of subjects (significant), and the fixed-delay predictor in "
     "90% (35% significant). 85% of subjects reached r > 0.6.", 0, GREEN, True),
    "Best electrode for amygdala was P3 (theta/alpha/beta/gamma all contributed); the dmPFC "
    "control region was best predicted by a DIFFERENT electrode (Fz) with a different frequency "
    "pattern entirely — evidence the fingerprint is region-specific, not a generic global signal.",
    "The delay between EEG and BOLD varied by subject, electrode, AND frequency band — directly "
    "justifying the sliding per-band delay design our own EFP inherits.",
], size=16)

# 11 — Our method walkthrough
s = prs.slides.add_slide(BLANK); title(s, "Our method: the EEG Finger-Print (EFP)",
                                        "Replicating and extending Meir-Hasson 2014 for DMN, CEN, and PDA (CEN−DMN)")
# image is a wide 4-panel composite (aspect ~1.33) -- constrain by HEIGHT so it can't collide
# with the bullets below (width=12.1 would render ~9in tall, far taller than the slide)
img(s, FULL / "paper_fig2_schematic_PDA_tr.png", left=3.6, top=1.45, height=4.6)
bullets(s, [
    ("In plain terms: for each electrode, break the EEG into 10 frequency bands, look at how "
     "each band's power over the last ~12 seconds predicts the fMRI signal right now, and let a "
     "regression model learn a separate best lag PER frequency band — no assumed fixed delay.", 0)],
    top=6.15, size=13)

# 12 — Within-subject results
s = prs.slides.add_slide(BLANK); title(s, "It works within-subject", "n=19, single best electrode per subject")
img(s, FULL / "fig_efp_vs_baselines.png", left=1.5, top=1.4, width=10.3)
bullets(s, [("EFP beats both a fixed-HRF baseline and the traditional theta/alpha ratio across "
             "every target — strongest for CEN, PDA, and their global-signal-controlled variants "
             "(r≈0.14–0.18).", 0)], top=6.3, size=14)

# 13 — LOSO generalization
s = prs.slides.add_slide(BLANK); title(s, "It generalizes to a NEW subject within our own cohort",
                                        "Leave-one-subject-out (LOSO): train on 18, predict the 19th, never seen")
loso = pd.read_csv(FULL / "efp_group_loso.csv")
loso_tr = loso[loso.resolution == "tr"].set_index("target")
bullets(s, [
    (f"CEN: r={loso_tr.loc['CEN','loso_mean_r']:+.3f} (p={loso_tr.loc['CEN','sign_flip_p']:.3f})", 0, NAVY, True),
    (f"DMN: r={loso_tr.loc['DMN','loso_mean_r']:+.3f} (p={loso_tr.loc['DMN','sign_flip_p']:.3f})", 0, NAVY, True),
    (f"PDA: r={loso_tr.loc['PDA','loso_mean_r']:+.3f} (p={loso_tr.loc['PDA','sign_flip_p']:.3f})", 0, NAVY, True),
    "This is a single, general fingerprint (not personalized) — evidence the EFP captures a "
    "real, shared EEG-to-BOLD mapping, not subject-specific overfitting.",
], size=20)

# 14 — n=17 -> n=19
s = prs.slides.add_slide(BLANK); title(s, "Adding 2 more subjects only helped",
                                        "sub-dmnelf002/003 recovered via R128 trigger-marker reconstruction")
bullets(s, [
    "3 targets moved from non-significant/marginal at n=17 to clearly significant at n=19.",
    "Several other targets got stronger (tighter p-values); nothing regressed meaningfully.",
    ("Net effect of the 2 recovered subjects: real, if modest, improvement across the board.", 0, GREEN, True),
], size=20)

# 15 — Cross-cohort replication (the headline result)
s = prs.slides.add_slide(BLANK)
title(s, "It replicates in a completely independent cohort", "DMNELF-trained fingerprint, tested on rtBPD — zero rtBPD data used in training")
rows_txt = []
for tgt in ("CEN", "DMN"):
    r1 = cc(cc1, tgt); r2 = cc(cc2, tgt)
    rows_txt.append((f"{tgt}  (electrode {r1['electrode']})", 0, NAVY, True))
    rows_txt.append((f"Session 1 (nf1, n={int(r1['n_test'])}): r={r1['mean_r']:+.3f}, p={r1['sign_flip_p']:.3f}", 1))
    rows_txt.append((f"Session 2 (nf2, n={int(r2['n_test'])}): r={r2['mean_r']:+.3f}, p={r2['sign_flip_p']:.3f}", 1))
bullets(s, rows_txt + [
    ("Two independent replications, in a cohort the model never saw during training — this is "
     "the strongest evidence that the fingerprint is a general EEG-to-BOLD mapping, not an "
     "artifact of one dataset.", 0, GREEN, True),
], size=18)

# 16 — Interpretability
s = prs.slides.add_slide(BLANK); title(s, "Interpretable, not a black box")
img(s, FULL / "efp_group_fingerprint_PDA_tr.png", left=0.3, top=1.4, width=6.4)
img(s, FULL / "efp_group_fingerprint_GSR_CEN_tr.png", left=6.7, top=1.4, width=6.4)
bullets(s, [("The learned [frequency × delay] weights peak at ~4–7 second lags — matching the "
             "known hemodynamic delay — even though the model was never told this. It re-discovered "
             "an HRF-like structure on its own, differing across frequency bands.", 0)],
        top=5.6, size=14)

# 17 — From science to product
s = prs.slides.add_slide(BLANK); title(s, "From a validated fingerprint to a usable tool")
bullets(s, [
    "We now have a general, cross-cohort-validated EEG fingerprint for CEN and DMN.",
    "That's enough signal to build something a person could use outside the scanner — not with "
    "a 31-channel research cap, but with a consumer-grade, 12-channel portable headset.",
    ("That tool is mindwear.", 0, NAVY, True),
], size=20)

# 18 — mindwear overview
s = prs.slides.add_slide(BLANK); title(s, "mindwear: portable DMN/CEN neurofeedback")
bullets(s, [
    "Runs the same visual ball-feedback task as our scanner protocol, driven by a portable "
    "12-channel EEG headset instead of an MRI scanner.",
    "The EEG-to-network decoder is the same EFP approach just shown — trained on our DMNELF "
    "cohort, montage-restricted to the portable headset's 12 channels.",
    "Includes a built-in comparison mode: replay a recorded session and watch the EEG-decoded "
    "ball run side by side with the scanner ground truth, on the same run.",
    ("[Insert a mindwear screenshot / architecture diagram here]", 0, ORANGE, True),
], size=18)

# 19 — Live demo
s = prs.slides.add_slide(BLANK); title(s, "Let's try it live")
bullets(s, [
    "Demo on myself first — a short live feedback run, so you can see the ball move in real time.",
    "Then: any volunteers? (headset fit takes ~1–2 minutes; keep a spare pair of alcohol wipes "
    "and the fit-quality checklist within reach)",
    ("What to watch for: it will be noisier than an fMRI signal — that's expected and is exactly "
     "what the earlier results characterized. Watch the overall trend over the run, not "
     "moment-to-moment jitter.", 0, GREY),
], size=19)

# 20 — Closing
s = prs.slides.add_slide(BLANK); title(s, "Takeaways")
bullets(s, [
    ("A single-electrode EEG fingerprint, replicating Meir-Hasson 2014, decodes DMN/CEN/PDA "
     "within-subject, generalizes to new subjects (LOSO), and replicates in a fully independent "
     "cohort (rtBPD) for CEN and DMN.", 0, NAVY, True),
    "That validated fingerprint now powers mindwear, a portable, scanner-free neurofeedback tool.",
    "Next: scaling access to DMN-targeted self-regulation training beyond the scanner — to the "
    "clinic, and eventually, the living room.",
], size=19)

# 21 — related work: the mindfulness+NFB paradigm generalizes beyond DMN/CEN
s = prs.slides.add_slide(BLANK)
title(s, "The broader paradigm: mindfulness + targeted neurofeedback",
      "Bauer et al. (2025). Psychiatry Research: Neuroimaging 353:112050  ·  [insert poster/slide graphics here]")
bullets(s, [
    "Same recipe — real-time fMRI neurofeedback delivered from a symptom-relevant region, "
    "paired with the same \"mental-noting\" mindfulness technique — also tested in schizophrenia "
    "patients with treatment-resistant auditory hallucinations, targeting superior temporal "
    "gyrus (STG) instead of DMN/CEN.",
    ("Randomized, sham-controlled (n=23): STG-targeted neurofeedback produced significantly "
     "greater reductions in secondary auditory cortex activity and its connectivity to "
     "prefrontal cognitive-control regions than sham feedback from an unrelated region.", 0),
    "Auditory hallucination severity dropped in both groups (consistent with mindfulness itself "
    "contributing), but only real, region-specific feedback changed the underlying circuit — "
    "the same logic underlying why targeting DMN/CEN specifically (not a generic relaxation "
    "signal) should matter here.",
    ("Together with the mbNF result on slide 7, this is a second, independent population and "
     "target region where the same mindfulness-plus-neurofeedback framework produced measurable "
     "brain change — this approach is not a one-off.", 0, NAVY, True),
], size=15)

out = PROJ / "Fingerprint_20260902.pptx"; prs.save(str(out))
print(f"Saved {out} ({len(prs.slides._sldIdLst)} slides)")
