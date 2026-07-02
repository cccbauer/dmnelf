#!/usr/bin/env python3
"""Build a summary PPTX for the dmn_hmm_detection (Cooray 2024 TIDE-HMM replication)."""
from pathlib import Path
import numpy as np, pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from scipy.stats import ttest_1samp
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJ = Path(__file__).resolve().parent.parent
RES = PROJ / "results" / "group_k12"
FIG = RES / "figures"
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xB0, 0x3A, 0x2E)

# ── Build feedback bar-chart figure ──
fb = pd.read_csv(RES / "feedback_state_correlations.csv")
targets = ["DMN", "CEN", "PDA", "GSR_DMN", "GSR_CEN", "GSR_PDA"]
means, sems, ps = [], [], []
for t in targets:
    v = fb[t].dropna().values
    means.append(v.mean()); sems.append(v.std(ddof=1) / np.sqrt(len(v)))
    ps.append(ttest_1samp(v, 0.0)[1])
fig, ax = plt.subplots(figsize=(8, 4.2))
colors = [GREEN.__str__() if p < 0.05 else "#999999" for p in ps]
xs = np.arange(len(targets))
ax.bar(xs, means, yerr=sems, color=["#2E7D32" if p < 0.05 else "#9e9e9e" for p in ps],
       capsize=4, edgecolor="k", linewidth=0.6)
ax.axhline(0, color="k", lw=0.8)
for i, (m, p) in enumerate(zip(means, ps)):
    star = "***" if p < 0.001 else "**" if p < 0.01 else "*" if p < 0.05 else "ns"
    ax.text(i, m + np.sign(m) * (sems[i] + 0.008), star, ha="center",
            va="bottom" if m >= 0 else "top", fontsize=11, fontweight="bold")
ax.set_xticks(xs); ax.set_xticklabels(targets, rotation=20)
ax.set_ylabel("Pearson r (State 7 occupancy vs fMRI)")
ax.set_title("Feedback runs (n=67): HMM State-7 occupancy vs fMRI network signal")
fig.tight_layout()
fb_fig = FIG / "feedback_state7_bars.png"
fig.savefig(fb_fig, dpi=140); plt.close(fig)

# rest correlations for the identification slide
rest = pd.read_csv(RES / "state_fmri_correlations.csv")


def r_for(df, target, state):
    row = df[(df.target == target) & (df.state == state)]
    return float(row.mean_r.iloc[0]) if len(row) else np.nan


# ── Presentation ──
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(15); r2.font.color.rgb = GREY
    return slide


def add_bullets(slide, bullets, top=1.5, left=0.6, width=12.1, height=5.5, size=17):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        lvl = 0; txt = b; color = None; bold = False
        if isinstance(b, tuple):
            txt, lvl = b[0], b[1]
            if len(b) > 2: color = b[2]
            if len(b) > 3: bold = b[3]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl);
        if color: r.font.color.rgb = color
        r.font.bold = bold
        p.space_after = Pt(6)


def add_img(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        kw = {}
        if width: kw["width"] = Inches(width)
        if height: kw["height"] = Inches(height)
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


# Slide 1 — title
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Real-Time DMN Detection with a Time-Delay Embedded HMM"
r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = NAVY
for line, sz, col in [
    ("Replicating Cooray et al. (2024) on the DMNELF simultaneous EEG–fMRI cohort", 20, GREY),
    ("K=12 TIDE-HMM · 17 subjects · validated against concurrent fMRI (their design could not)", 15, GREY)]:
    p = tf.add_paragraph(); rr = p.add_run(); rr.text = line
    rr.font.size = Pt(sz); rr.font.color.rgb = col

# Slide 2 — background & question
s = prs.slides.add_slide(BLANK)
add_title(s, "Background & question")
add_bullets(s, [
    "Cooray et al. 2024 used a Time-Delay Embedded HMM (osl-dynamics) to detect a DMN state from EEG in real time (LEMON dataset).",
    ("They had NO simultaneous fMRI — the DMN state was identified purely by its spectral signature (posterior alpha + frontal delta/theta).", 1),
    "DMNELF advantage: simultaneous EEG–fMRI lets us validate the HMM's DMN state DIRECTLY against concurrent fMRI DMN — a stronger test.",
    "Questions:",
    ("1. Does a DMN-like HMM state exist, and does it track real fMRI DMN?", 1),
    ("2. Does the spectral heuristic pick the same state as fMRI validation?", 1),
    ("3. Is single-state occupancy a useful real-time decoder vs our DWT+stats regression?", 1),
])

# Slide 3 — methods
s = prs.slides.add_slide(BLANK)
add_title(s, "Methods")
add_bullets(s, [
    "EEG: resting-state, 1–45 Hz, 250 Hz; 31 channels; 17 subjects × 2 runs (33 rest runs, ~3.5M samples).",
    "Feature prep: time-delay embedding (n=7) + PCA→80 components (99.8% variance), whitened & standardized.",
    "Model: HMM, K=12 states, learn_covariances (osl-dynamics); random-subset initialization ×10, 20 epochs.",
    "DMN-state identification (our contribution): bin each state's occupancy to TR (1.2 s), HRF-convolve, correlate vs simultaneous fMRI DMN / CEN / PDA (+ GSR'd).",
    "Cross-check: per-state multitaper spectra → posterior-alpha / frontal-theta signature (paper's method).",
    "Feedback validation: apply trained HMM (same PCA transform) to 67 feedback runs; correlate State-7 occupancy vs fMRI; compare to regression benchmark.",
], size=16)

# Slide 4 — model health
s = prs.slides.add_slide(BLANK)
add_title(s, "The HMM fits cleanly — 12 healthy states")
add_bullets(s, [
    "All 12 states used (fractional occupancy 0.03–0.12); 100% of timepoints confidently assigned (max-alpha > 0.5).",
    "Distinct state covariances (pairwise Frobenius 5.8–83); transition stickiness 0.50–0.97.",
    ("State 12 = high-amplitude artifact/movement state (covariance trace 694 vs ~30–77) — excluded from interpretation.", 1),
    ("Note: an early 'degenerate occupancy' printout was a reporting bug (one-hot vs integer), not model collapse.", 1, GREY),
], height=2.2)
add_img(s, FIG / "state_topomaps_alpha.png", left=2.4, top=3.0, width=8.5)

# Slide 5 — DMN state identification (rest)
s = prs.slides.add_slide(BLANK)
add_title(s, "State 7 = DMN state in rest — but weak")
add_bullets(s, [
    f"State 7 is the DMN state across all three fMRI criteria (rest, 33 runs):",
    (f"+corr with raw fMRI DMN: r = {r_for(rest,'DMN',7):+.3f}", 1),
    (f"+corr with GSR'd DMN:    r = {r_for(rest,'GSR_DMN',7):+.3f}", 1),
    (f"−corr with PDA (CEN−DMN): r = {r_for(rest,'PDA',7):+.3f}", 1),
    "Internally consistent (DMN-positive, PDA-negative) → a real DMN-like state exists.",
    ("BUT weak: r ≈ 0.05–0.06, vs multivariate DWT+stats regression GSR_DMN r≈0.12 / GSR_CEN r≈0.17.", 0, RED),
    "A single state's occupancy is a weaker DMN readout than continuous multivariate regression.",
])

# Slide 6 — spectral disagreement
s = prs.slides.add_slide(BLANK)
add_title(s, "The spectral heuristic disagrees with fMRI")
add_bullets(s, [
    "Paper's DMN signature = posterior ALPHA high + frontal DELTA/THETA high.",
    ("State 7 (the fMRI-DMN state) has LOW alpha: frontal rank 11/12, posterior rank 9/12.", 1, RED),
    "The classic posterior-alpha state (State 10) is slightly NEGATIVE with fMRI DMN (r≈−0.05) — consistent with the known alpha↔DMN-BOLD anticorrelation.",
    "→ In our data the spectral heuristic would pick the WRONG state.",
    ("Only detectable because we have simultaneous fMRI — the caveat the original design could not see.", 0, NAVY, True),
], height=2.4)
add_img(s, FIG / "dmn_state_spectrum.png", left=3.3, top=3.4, width=6.7)

# Slide 7 — feedback validation
s = prs.slides.add_slide(BLANK)
add_title(s, "Feedback: the DMN label shifts — State 7 tracks CEN/PDA")
add_img(s, fb_fig, left=0.4, top=1.5, width=7.4)
tb = s.shapes.add_textbox(Inches(8.0), Inches(1.6), Inches(5.0), Inches(5.4))
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, col, bold) in enumerate([
    ("Across 67 feedback runs, State-7 occupancy:", None, True),
    ("GSR_CEN  r = +0.115  (p<0.001)", GREEN, True),
    ("GSR_PDA  r = +0.096  (p=0.001)", GREEN, True),
    ("PDA        r = +0.094  (p=0.001)", GREEN, True),
    ("DMN / GSR_DMN ≈ null-to-negative", RED, False),
    ("", None, False),
    ("The state that was DMN-positive in REST tracks CEN/PDA in FEEDBACK — its identity is task-dependent.", NAVY, True),
    ("", None, False),
    ("PDA (CEN↑/DMN↓) is the neurofeedback target: a single interpretable, real-time-friendly state decodes it significantly.", None, False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = txt; r.font.size = Pt(15)
    if col: r.font.color.rgb = col
    r.font.bold = bold; p.space_after = Pt(6)

# Slide 8 — verdict
s = prs.slides.add_slide(BLANK)
add_title(s, "Verdict")
add_bullets(s, [
    ("What worked", 0, NAVY, True),
    ("TIDE-HMM fits cleanly on 17 subjects; a DMN-like state exists and was validated against concurrent fMRI.", 1),
    ("In feedback, State-7 occupancy significantly decodes GSR_CEN (r=0.115***) and PDA (r=0.096**) — the neurofeedback target.", 1),
    ("What simultaneous fMRI revealed (novel vs the paper)", 0, NAVY, True),
    ("The spectral heuristic picks the wrong state — the fMRI-DMN state is not the posterior-alpha one.", 1),
    ("The 'DMN state' identity is task-dependent (rest→DMN-ish, feedback→CEN/PDA).", 1),
    ("Recommendation", 0, NAVY, True),
    ("DWT+stats multivariate regression remains the primary decoder (stronger, task-stable).", 1),
    ("HMM is a valuable complementary / interpretability tool and a candidate real-time single-state readout, with the task-shift caveat.", 1),
], size=16)

out = PROJ / "dmn_hmm_detection_results.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
