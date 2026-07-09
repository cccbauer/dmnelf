#!/usr/bin/env python3
"""Build the EEG f-SNR two-stream deck (pedagogical style, matches fsnr/efp decks)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PROJ = Path(__file__).resolve().parent.parent
R = PROJ / "results"
NAVY = RGBColor(0x1F, 0x3A, 0x5F); GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32); RED = RGBColor(0xB0, 0x3A, 0x2E); BLUE = RGBColor(0x1F, 0x77, 0xB4)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def title(s, t, sub=None):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)); tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = t; r.font.size = Pt(27); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p = tf.add_paragraph(); rr = p.add_run(); rr.text = sub; rr.font.size = Pt(14); rr.font.color.rgb = GREY


def bullets(s, items, top=1.5, left=0.6, width=12.1, height=5.5, size=17):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)); tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(items):
        lvl, txt, col, bold = 0, b, None, False
        if isinstance(b, tuple):
            txt, lvl = b[0], b[1]; col = b[2] if len(b) > 2 else None; bold = b[3] if len(b) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl)
        if col: r.font.color.rgb = col
        r.font.bold = bold; p.space_after = Pt(6)


def img(s, path, left, top, width=None):
    if Path(path).exists():
        kw = {"width": Inches(width)} if width else {}
        s.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def new(t=None, sub=None):
    s = prs.slides.add_slide(BLANK)
    if t: title(s, t, sub)
    return s


# 1 title
s = new()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.0)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "An EEG f-SNR that matches the BOLD neurofeedback target"
r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
for line, sz in [("Two streams: (A) fit an EEG decoder to the fMRI f-SNR/PDA  vs  (B) a PURE EEG f-SNR matched to it", 17),
                 ("The prize: a calibration-free, interpretable, portable EEG 'clarity' index that tracks PDA", 15),
                 ("DMNELF (n=17) within/LOSO + rtBPD (nf1 n=19, nf2 n=11) cross-cohort", 14)]:
    p = tf.add_paragraph(); rr = p.add_run(); rr.text = line; rr.font.size = Pt(sz); rr.font.color.rgb = GREY

# 2 idea
s = new("The idea — align f-SNR across modalities, don't just decode")
bullets(s, [
    ("fMRI phase: the f-SNR ≈ a normalized restatement of PDA regulation (signal channel).", 0),
    ("Instead of only fitting EEG→BOLD, compute an f-SNR in EACH modality and MATCH them — "
     "signal/noise 'clarity' in EEG vs in BOLD.", 0, NAVY, True),
    ("Stream A (benchmark): EFP / band-power decoder fit to the fMRI f-SNR / PDA (the ceiling).", 0),
    ("Stream B (the prize): a PURE EEG f-SNR (no per-subject fitting) — if it matches PDA and "
     "TRANSFERS, we get a calibration-free NF marker.", 0, BLUE, True),
    ("EEG f-SNR flavors searched: band power (running mean/std + quench) and oscillatory÷aperiodic "
     "(specparam 1/f).", 0),
], top=1.6, size=17)

# 3 within head-to-head
s = new("Within-subject head-to-head — construct reaches ~70% of the fitted ceiling",
        "Leak-free match to PDA. Fitted decoder leads; the zero-fitting construct f-SNR is close behind.")
img(s, R / "fig_headtohead_within.png", left=1.9, top=1.9, width=9.5)
bullets(s, [("Fitted EFP ~0.17; construct frontal-theta f-SNR 0.119 (NO fitting). specparam osc/aperiodic "
             "is weaker (0.10) — per-TR FOOOF adds noise, simpler band power wins.", 0)], top=6.4, size=12.5)

# 4 the payoff — cross-cohort
s = new("The payoff — the pure EEG f-SNR TRANSFERS across cohorts; raw power does not",
        "Fixed frontal-theta f-SNR construct, zero fitting, applied to an independent cohort (rtBPD)")
img(s, R / "fig_crosscohort_generalize.png", left=1.9, top=1.9, width=9.5)
bullets(s, [("rtBPD nf1 r=+0.126 (p=5e-4, 79% subj+), nf2 r=+0.147 (p=5e-3, 91%+). Raw band power collapses "
             "(r≈0, ns) — noise-normalization removes cohort-specific gain. Matches/beats the fitted EFP.", 0, GREEN, True)],
        top=6.4, size=12.5)

# 5 why it works
s = new("Why the f-SNR transfers when raw EEG doesn't")
bullets(s, [
    ("Raw band power carries the PDA signal (|r|≈0.12 within-cohort) but its scale/gain is "
     "cohort-specific → r≈0 cross-cohort.", 0),
    ("The f-SNR = signal/noise DIVIDES OUT that gain → a dimensionless 'clarity' that is "
     "cohort-invariant. This is exactly the framework's claim, borne out.", 0, NAVY, True),
    ("Frontal montage works (portable headset); frontal-theta is physiologically apt for "
     "DMN↔CEN control. The f-SNR normalization specifically helps the noisier frontal channels.", 0),
    ("A zero-fitting construct that matches/beats a per-cohort-fitted decoder cross-cohort is the "
     "ideal NF marker: interpretable, calibration-free, deployable.", 0, GREEN, True),
], top=1.6, size=16)

# 6 secondary — EEG quench
s = new("Secondary — clean EEG variability quenching")
bullets(s, [
    ("Cross-modal analog of the BOLD result: does EEG variance declutter during feedback?", 0),
    ("Non-convolved (specparam) band power: beta +1.6 dB (p=2e-4), gamma +3.7 dB (p=1e-4) variance "
     "DROP during feedback; delta/theta/alpha flat.", 0, BLUE, True),
    ("High-frequency EEG decluttering during regulation — but caveat: beta/gamma are EMG-sensitive "
     "(could be the subject sitting stiller).", 0),
    ("(The HRF-convolved cache gave a spurious +30 dB in all bands — onset-ramp artifact; corrected "
     "with the non-convolved specparam extraction.)", 0, GREY),
], top=1.7, size=16)

# 7 bottom line
s = new("Bottom line")
bullets(s, [
    ("A PURE EEG f-SNR (frontal-theta, band power, zero fitting) matches the BOLD PDA the "
     "neurofeedback trains — within-subject 0.12, and it TRANSFERS to an independent cohort "
     "(rtBPD nf1 0.126, nf2 0.147).", 0, NAVY, True),
    ("It generalizes where raw EEG and even the fitted decoder struggle cross-cohort, because "
     "noise-normalization removes cohort-specific gain.", 0, GREEN),
    ("Practical NF implication: a calibration-free, portable-frontal EEG clarity index — no "
     "per-subject/per-cohort training — can serve as an EEG-only proxy for the fMRI NF target.", 0, GREEN, True),
    ("The fitted EFP decoder still wins within-subject (0.17) — use it when calibration is "
     "available; use the construct f-SNR when it isn't (new subjects/sites).", 0),
], top=1.7, size=16)

out = PROJ / "eeg_fsnr_results.pptx"; prs.save(str(out))
print(f"Saved {out} ({len(prs.slides._sldIdLst)} slides)")
