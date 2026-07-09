#!/usr/bin/env python3
"""Build the fMRI f-SNR investigation deck (pedagogical style, matches efp/bold decks)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PROJ = Path(__file__).resolve().parent.parent
R = PROJ / "results"
NAVY = RGBColor(0x1F, 0x3A, 0x5F); GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32); RED = RGBColor(0xB0, 0x3A, 0x2E); BLUE = RGBColor(0x1F, 0x77, 0xB4)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def title(s, text, sub=None):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0)); tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text; r.font.size = Pt(27); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p = tf.add_paragraph(); rr = p.add_run(); rr.text = sub; rr.font.size = Pt(14); rr.font.color.rgb = GREY


def bullets(s, items, top=1.5, left=0.6, width=12.1, height=5.5, size=17):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)); tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(items):
        lvl, txt, col, bold = 0, b, None, False
        if isinstance(b, tuple):
            txt, lvl = b[0], b[1]; col = b[2] if len(b) > 2 else None; bold = b[3] if len(b) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - 2 * lvl)
        if col: r.font.color.rgb = col
        r.font.bold = bold; p.space_after = Pt(6)


def img(s, path, left, top, width=None, height=None):
    if Path(path).exists():
        kw = {}
        if width: kw["width"] = Inches(width)
        if height: kw["height"] = Inches(height)
        s.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


def new(t=None, sub=None):
    s = prs.slides.add_slide(BLANK)
    if t: title(s, t, sub)
    return s


# 1 — title
s = new()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.0)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = "A functional signal-to-noise ratio (f-SNR) for DMNELF neurofeedback"
r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
for line, sz in [("Purely-fMRI phase — applying the f-SNR framework (Laukkonen 2026; Nath 2026) to CEN / DMN / PDA", 18),
                 ("Is there an fMRI f-SNR? Does increased PDA / reduced DMN = higher f-SNR? Is it a good NF signal?", 15),
                 ("17 subjects · 67 feedback runs · law of total variance · before the EEG decoder", 14)]:
    p = tf.add_paragraph(); rr = p.add_run(); rr.text = line; rr.font.size = Pt(sz); rr.font.color.rgb = GREY

# 2 — the framework
s = new("The f-SNR framework (the idea we're testing)")
bullets(s, [
    ("f-SNR = signal variance / noise variance — 'mental clarity' = neural activity that tracks "
     "the goal-relevant cause, not endogenous noise.", 0),
    ("Law of total variance:  Var(r) = Var_z(E[r|z])  [signal]  +  E_z(Var(r|z))  [noise].", 0, NAVY, True),
    ("Two channels: amplify signal (task-positive up) and declutter noise. Nath 2026 names "
     "DMN suppression as the noise-reduction substrate.", 0),
    ("Our hypothesis for DMNELF: increased PDA (CEN−DMN) / reduced DMN = higher f-SNR — testable, "
     "non-circular, because f-SNR is a variance RATIO while PDA/DMN are means.", 0, BLUE, True),
    ("Goal here: build a faithful fMRI f-SNR from the networks, then ask if it is a good "
     "neurofeedback target — before hunting its EEG signature.", 0),
], top=1.6, size=17)

# 3 — method
s = new("Method — law of total variance, within-run rest vs feedback")
bullets(s, [
    ("Every feedback run = 25 TR rest baseline + 100 TR feedback (drop 5 HRF-lag TRs) → a "
     "controlled within-run cause z = {rest, feedback}.", 0),
    ("signal = Var_z(E[r|z]) (how far the network mean shifts rest→feedback);  "
     "noise = E_z(Var(r|z)) (within-condition fluctuation).", 0),
    ("f-SNR = signal / noise, in dB. A GLM version uses the HRF-convolved boxcar (the "
     "pseudo-target) as E[r|z]  →  f-SNR = R²/(1−R²).", 0),
    ("Computed per network: PDA, CEN, DMN. n = 17 subjects, 67 runs.", 0, GREY),
], top=1.7, size=17)

# 4 — result 1
s = new("Result 1 — f-SNR tracks the SIGNAL, not DMN mean-suppression",
        "Group regulates correctly (PDA+ 78%, CEN+ 79%, DMN− 60%). f-SNR ↔ PDA regulation (r=+0.58); f-SNR ↔ DMN mean drop ≈ 0")
img(s, R / "fig_fsnr_vs_pda_dmn.png", left=0.5, top=1.85, width=12.3)
bullets(s, [("'Increased PDA → higher f-SNR' holds; 'reduced DMN → higher f-SNR' does NOT, as a mean effect.", 0, NAVY, True)],
        top=6.4, size=13)

# 5 — reframe to variance
s = new("Reframe — DMN is the NOISE, so test its VARIANCE (variability quenching)",
        "Nath Table 2 row 1: decluttering = reduced endogenous variance during regulation")
img(s, R / "fig_dmn_quench.png", left=0.5, top=1.85, width=12.3)
bullets(s, [("DMN variance QUENCHES rest→feedback (+2.2 dB, p=1e-4, 75% of runs), more than CEN — the "
             "noise-reduction the mean test missed. But the degree of quench does not predict regulation.", 0, GREEN, True)],
        top=6.4, size=13)

# 6 — tighten
s = new("Tightening — real, but GLOBAL not DMN-specific, and dissociable from signal",
        "Startup control passes (long rest = no quench); raw global_signal as whole-brain reference")
img(s, R / "fig_fsnr_tighten.png", left=0.5, top=1.85, width=12.3)
bullets(s, [
    ("Control passes: dedicated long-rest shows NO quench → feedback quench is real, not an onset artifact.", 0, GREEN),
    ("But global quenches MOST (+3.1) > DMN (+2.2) > CEN (+1.25); and DMN quench is INDEPENDENT of f-SNR / "
     "regulation (r≈0). Noise-down and signal-up are dissociable axes.", 0, RED, True),
], top=6.15, size=12.5)

# 7 — timeseries
s = new("Seeing it — the timeseries", "Group-mean trajectory · quench trajectory · rest control · running f-SNR(t) in dB")
img(s, R / "fig_fsnr_timeseries_group.png", left=0.15, top=2.0, width=13.0)
bullets(s, [("After feedback onset: PDA/CEN rise (signal), variance drops (global most), and running f-SNR(t) rises "
             "to a peak ~+3.5 dB then decays. DMN f-SNR lags. (Per-run browser: fsnr_timeseries_runs.pdf, 67 pages.)", 0)],
        top=6.5, size=12.5)

# 8 — NF proxy verdict
s = new("Is f-SNR a good neurofeedback signal? — valid, but not better than PDA",
        "Causal real-time running f-SNR = trailing_mean(PDA)/trailing_std")
img(s, R / "fig_fsnr_proxy.png", left=0.7, top=1.9, width=11.9)
bullets(s, [
    ("Well-modulated (fb>rest +1.07, p=1e-4, 84%), smooth (0.95), reliable (ICC 0.51), controllable (r=0.77 with β_PDA).", 0),
    ("BUT rest↔feedback discriminability d′ = 0.70 (f-SNR) < 0.82 (raw PDA) → normalizing by noise does NOT sharpen "
     "the signal. Raw PDA stays the better target.", 0, RED, True),
], top=6.1, size=12.5)

# 9 — bottom line
s = new("Bottom line")
bullets(s, [
    ("DMNELF neurofeedback engages BOTH f-SNR channels — but they are dissociable.", 0, NAVY, True),
    ("Signal channel (CEN/PDA regulation): carries the reliable, state-separating, controllable information.", 1, GREEN),
    ("Noise channel (variability quenching): real, but mostly GLOBAL, and independent of the signal / of NF success.", 1),
    ("f-SNR is a valid interpretive index of 'clarity' but NOT a better NF target than raw PDA "
     "(noise-normalization doesn't help here).", 0, RED, True),
    ("For the EEG phase we carry the SIGNAL-channel target (glm_PDA_db / PDA); the global decluttering "
     "axis is a separate arousal-like signal EEG may read but that doesn't serve the NF goal.", 0, NAVY, True),
    ("Faithful to the framework, honestly tested: the 'clear mind = higher f-SNR' we can use is the "
     "amplify-signal channel, not declutter-DMN-noise — in this paradigm.", 0),
], top=1.6, size=16)

# 10 — next
s = new("Next — the EEG decoder")
bullets(s, [
    ("Target: the signal-channel fMRI f-SNR / PDA (reliable, controllable).", 0),
    ("Reuse the EFP (single-electrode sliding-delay) and multivariate band-power machinery already built.", 0),
    ("Ask: can EEG decode the fMRI f-SNR / PDA, and does the pseudo-target calibration transfer?", 0, BLUE, True),
    ("Secondary: is the GLOBAL decluttering axis the arousal signal EEG reads best (band-power's known confound)?", 0),
], top=1.7, size=17)

out = PROJ / "fsnr_results.pptx"; prs.save(str(out))
print(f"Saved {out} ({len(prs.slides._sldIdLst)} slides)")
